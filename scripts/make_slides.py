"""Build the capstone defence deck from the results of record.

Produces ``slides/MSE_Capstone_Defence.pptx`` -- 23 content slides plus 10
backup slides, sized for a 20-minute talk, with speaker notes on every slide.

Every figure quoted here is traced to ``docs/CHANGELOG.md``, which is the single
source of truth. Two figures are rendered fresh by this script rather than reused
from ``MSE_CAPSTONE_REPORT_new/figures/``, because the committed
``fusion_approaches.png`` predates the 2026-08-10 multi-seed repetition and still
carries the withdrawn single-seed 69.55% figure.

Run from the project root::

    python scripts/make_slides.py
"""

from __future__ import annotations

import shutil
import sys
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
SLIDE_DIR = ROOT / "slides"
FIG_DIR = SLIDE_DIR / "figures"
REPORT_FIGS = ROOT / "MSE_CAPSTONE_REPORT_new" / "figures"

# --- Palette -----------------------------------------------------------------
# Same validated palette as scripts/generate_result_figures.py, so the two new
# figures sit beside the report figures without a visible seam. Slot order is the
# colour-vision-deficiency safety mechanism and is not reshuffled.
SERIES_1 = "#2a78d6"   # blue   -- the recommended / zero-parameter systems
SERIES_2 = "#eb6834"   # orange -- single-model or learned-fusion systems
SERIES_3 = "#199e70"   # green  -- fusion series

INK_PRIMARY = "#0b0b0b"
INK_SECONDARY = "#52514e"
INK_MUTED = "#898781"
GRIDLINE = "#e1e0d9"
BASELINE = "#c3c2b7"
SURFACE = "#ffffff"
ACCENT = SERIES_1
WARN = "#b3401f"

FONT = "Segoe UI"
MONO = "Consolas"

SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.85)
CONTENT_W = SW - 2 * MARGIN


def rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str.lstrip("#").upper())


# =============================================================================
# Figures rendered fresh for the deck
# =============================================================================

def style_axes(ax, *, xgrid=False, ygrid=True):
    ax.set_facecolor(SURFACE)
    for side in ("top", "right"):
        ax.spines[side].set_visible(False)
    for side in ("left", "bottom"):
        ax.spines[side].set_color(BASELINE)
        ax.spines[side].set_linewidth(1.0)
    ax.tick_params(colors=INK_MUTED, labelsize=11, length=0)
    if ygrid:
        ax.yaxis.grid(True, color=GRIDLINE, linewidth=1.0)
    if xgrid:
        ax.xaxis.grid(True, color=GRIDLINE, linewidth=1.0)
    ax.set_axisbelow(True)


def fig_combiner_freedom(path: Path) -> None:
    """Each combiner's paired advantage over the audio model from its own run.

    Absolute accuracies are not comparable across these two runs -- their audio
    baselines differ by 2.19%, which is inside the noise floor -- so plotting
    the *paired difference against the same run's audio* is the only comparison
    that means anything. It also removes the withdrawn 69.55% absolute from the
    slide entirely, and shows that the ordering replicates across two
    independent runs rather than resting on one.

    Sources: outputs/late_fusion_cv.json (standard EEG encoder) and
    outputs/late_fusion_cv_final.json (adversarial EEG encoder).
    """
    rules = [
        ("Mean of probabilities", "0 params", +5.19, +5.07),
        ("Weighted average", "1 param", +4.26, +4.83),
        ("Per-class weights", "5 params", +4.69, +4.43),
        ("Max-confidence gating", "0 params, hard switch", +2.76, +3.24),
        ("Trained attention fusion", "8,064 params", +1.83, -0.36),
    ]
    rules = rules[::-1]
    labels = [f"{name}\n{p}" for name, p, _, _ in rules]

    fig, ax = plt.subplots(figsize=(12.4, 5.8), dpi=200)
    y = np.arange(len(rules))
    h = 0.34

    std = [r[2] for r in rules]
    adv = [r[3] for r in rules]
    ax.barh(y + h / 2 + 0.02, std, height=h, color=SERIES_1, label="Run A — standard EEG encoder")
    ax.barh(y - h / 2 - 0.02, adv, height=h, color=SERIES_3, label="Run B — adversarial EEG encoder")

    for yi, (a, b) in enumerate(zip(std, adv)):
        for val, off in ((a, h / 2 + 0.02), (b, -h / 2 - 0.02)):
            ha, dx = ("left", 0.09) if val >= 0 else ("right", -0.09)
            ax.text(val + dx, yi + off, f"{val:+.2f}", va="center", ha=ha,
                    fontsize=11, color=INK_PRIMARY, fontweight="bold")

    ax.axvline(0, color=INK_SECONDARY, linewidth=1.6, zorder=4)
    ax.set_yticks(y)
    ax.set_yticklabels(labels, fontsize=11.5, color=INK_SECONDARY)
    ax.set_ylim(-0.6, len(rules) - 0.4)
    ax.set_xlim(-1.6, 6.4)
    ax.set_xlabel("Advantage over the audio model from the same run, in % "
                  "(paired, 7-fold subject-wise CV)",
                  fontsize=12, color=INK_SECONDARY, labelpad=10)
    ax.xaxis.set_major_formatter(lambda v, _: f"{v:+.0f}")
    style_axes(ax, xgrid=True, ygrid=False)

    ax.axvspan(-1.0, 1.0, color=GRIDLINE, alpha=0.55, zorder=0)
    ax.text(0.0, len(rules) - 0.52, "noise floor, ±1%", ha="center", va="bottom",
            fontsize=10.5, color=INK_MUTED)

    ax.legend(loc="lower right", frameon=False, fontsize=11.5, labelcolor=INK_SECONDARY,
              handlelength=1.2, handleheight=1.0)

    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor=SURFACE)
    plt.close(fig)


def fig_multiseed(path: Path) -> None:
    """Per-seed paired comparison: audio vs late fusion, three repetitions.

    Two series, both directly labelled and in the legend, so identity survives
    greyscale and colour-vision deficiency.
    """
    audio = [61.55, 62.29, 64.48]
    fusion = [66.12, 66.69, 68.26]
    runs = ["Run 1", "Run 2", "Run 3"]

    fig, ax = plt.subplots(figsize=(11.0, 5.4), dpi=200)
    x = np.arange(len(runs))

    for xi, a, f in zip(x, audio, fusion):
        ax.plot([xi, xi], [a, f], color=BASELINE, linewidth=2.0, zorder=1)
        ax.annotate(f"+{f - a:.2f} pts", xy=(xi + 0.09, (a + f) / 2), fontsize=12,
                    color=INK_PRIMARY, va="center", ha="left", fontweight="bold")

    ax.plot(x, audio, "o-", color=SERIES_2, linewidth=2.0, markersize=11,
            markeredgecolor=SURFACE, markeredgewidth=2, label="Audio only", zorder=2)
    ax.plot(x, fusion, "o-", color=SERIES_1, linewidth=2.0, markersize=11,
            markeredgecolor=SURFACE, markeredgewidth=2,
            label="Late fusion (mean of probabilities)", zorder=2)

    for xi, a, f in zip(x, audio, fusion):
        ax.text(xi, a - 0.85, f"{a:.2f}%", ha="center", fontsize=11.5, color=INK_SECONDARY)
        ax.text(xi, f + 0.55, f"{f:.2f}%", ha="center", fontsize=11.5, color=INK_SECONDARY)

    ax.set_xticks(x)
    ax.set_xticklabels(runs, fontsize=12.5, color=INK_SECONDARY)
    ax.set_xlim(-0.45, len(runs) - 0.25)
    ax.set_ylim(58.5, 71.5)
    ax.set_ylabel("Accuracy, 7-fold subject-wise CV", fontsize=12, color=INK_SECONDARY,
                  labelpad=10)
    ax.yaxis.set_major_formatter(lambda v, _: f"{v:.0f}%")
    style_axes(ax)
    ax.legend(loc="upper left", frameon=False, fontsize=12, labelcolor=INK_SECONDARY)

    ax.text(0.5, -0.16,
            "Each run redraws both the fold assignment and the initialisation. "
            "Paired difference: +4.25% ± 0.41 (p = 0.0031).",
            transform=ax.transAxes, ha="center", fontsize=11.5, color=INK_MUTED)

    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor=SURFACE)
    plt.close(fig)


def _rbox(ax, x, y, w, h, label, sub=None, *, fill, ink=SURFACE, fs=11.5):
    ax.add_patch(plt.Rectangle((x, y), w, h, facecolor=fill, edgecolor="none",
                               zorder=2))
    cy = y + h / 2 + (0.10 if sub else 0)
    ax.text(x + w / 2, cy, label, ha="center", va="center", fontsize=fs,
            color=ink, fontweight="bold", zorder=3)
    if sub:
        ax.text(x + w / 2, y + h / 2 - 0.20, sub, ha="center", va="center",
                fontsize=fs - 2.0, color=ink, zorder=3)


def _arrow(ax, x1, y1, x2, y2, *, colour=None, style="->", lw=1.8):
    ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                arrowprops=dict(arrowstyle=style, color=colour or BASELINE,
                                linewidth=lw, shrinkA=0, shrinkB=0))


def _elbow(ax, points, *, colour=None, lw=1.8):
    """Orthogonal connector through waypoints, arrowhead on the final segment.

    Long diagonal connectors in a two-path diagram cross each other ambiguously;
    right-angled runs in separate lanes cross only at right angles, which reads
    as wiring rather than as a mistake.
    """
    colour = colour or BASELINE
    xs = [p[0] for p in points]
    ys = [p[1] for p in points]
    ax.plot(xs[:-1], ys[:-1], color=colour, linewidth=lw, solid_capstyle="round",
            solid_joinstyle="round", zorder=1)
    _arrow(ax, xs[-2], ys[-2], xs[-1], ys[-1], colour=colour, lw=lw)


def fig_architecture(path: Path) -> None:
    """The system: one shared front end, two competing ways to combine it.

    Drawn as a single figure because the comparison between the two paths is the
    thesis's central experiment, and separating them into two figures would hide
    that they share identical encoders.
    """
    fig, ax = plt.subplots(figsize=(13.0, 6.2), dpi=200)
    ax.set_xlim(0, 13)
    ax.set_ylim(0, 6.2)
    ax.axis("off")

    # --- shared front end ---------------------------------------------------
    _rbox(ax, 0.15, 4.05, 1.95, 0.95, "EEG signal", "30 channels × 2,500", fill="#6da7ec")
    _rbox(ax, 0.15, 1.25, 1.95, 0.95, "Speech audio", "64 mel bands × 1,313", fill="#f2a07a")

    _rbox(ax, 2.55, 4.05, 2.05, 0.95, "EEG encoder", "band power → 128-d", fill=SERIES_1)
    _rbox(ax, 2.55, 1.25, 2.05, 0.95, "Audio encoder", "3 conv blocks → 128-d", fill=SERIES_2)

    _arrow(ax, 2.12, 4.52, 2.52, 4.52)
    _arrow(ax, 2.12, 1.72, 2.52, 1.72)

    ax.text(0.15, 5.42, "Shared front end — identical in both paths",
            ha="left", fontsize=11.5, color=INK_SECONDARY, style="italic")

    # --- path A: learned fusion --------------------------------------------
    ax.add_patch(plt.Rectangle((5.15, 3.35), 7.6, 2.35, facecolor="#f2f6fb",
                               edgecolor="none", zorder=0))
    ax.text(5.35, 5.42, "PATH A   Cross-modal attention fusion", fontsize=12,
            color=SERIES_1, fontweight="bold")
    ax.text(12.55, 5.42, "8,064 fusion parameters", fontsize=11,
            color=INK_SECONDARY, ha="right")

    _rbox(ax, 5.45, 3.85, 2.75, 1.15, "Cross-modal attention",
          "each signal attends to the other", fill=SERIES_1, fs=11)
    _rbox(ax, 8.55, 3.85, 1.85, 1.15, "Gated blend", "learned weight g", fill=SERIES_1,
          fs=11)
    _rbox(ax, 10.75, 3.85, 1.85, 1.15, "Classifier", "128-d → 5", fill="#256abf", fs=11)

    # EEG into path A: straight across its own lane.
    _elbow(ax, [(4.62, 4.62), (5.42, 4.62)], colour="#6da7ec")
    # Audio into path A: rises in the x = 5.05 lane.
    _elbow(ax, [(4.62, 1.80), (5.05, 1.80), (5.05, 4.22), (5.42, 4.22)],
           colour="#f2a07a")
    _arrow(ax, 8.22, 4.42, 8.52, 4.42)
    _arrow(ax, 10.42, 4.42, 10.72, 4.42)

    # --- path B: late fusion ------------------------------------------------
    ax.add_patch(plt.Rectangle((5.15, 0.45), 7.6, 2.35, facecolor="#eef7f2",
                               edgecolor="none", zorder=0))
    ax.text(5.35, 2.52, "PATH B   Late fusion — average the two answers",
            fontsize=12, color=SERIES_3, fontweight="bold")
    ax.text(12.55, 2.52, "0 fusion parameters", fontsize=11, color=INK_SECONDARY,
            ha="right")

    _rbox(ax, 5.45, 1.55, 2.15, 0.72, "EEG classifier", fill=SERIES_1, fs=10.5)
    _rbox(ax, 5.45, 0.68, 2.15, 0.72, "Audio classifier", fill=SERIES_2, fs=10.5)
    _rbox(ax, 8.05, 1.05, 2.35, 1.15, "Mean of the two", "probability outputs",
          fill=SERIES_3, fs=11)
    _rbox(ax, 10.75, 1.05, 1.85, 1.15, "Prediction", "5 emotions", fill="#199e70", fs=11)

    # EEG into path B: descends in the x = 4.80 lane, clear of the audio lane.
    _elbow(ax, [(4.62, 4.42), (4.80, 4.42), (4.80, 1.91), (5.42, 1.91)],
           colour="#6da7ec")
    # Audio into path B: descends in its own lane.
    _elbow(ax, [(4.62, 1.64), (5.20, 1.64), (5.20, 1.04), (5.42, 1.04)],
           colour="#f2a07a")
    _arrow(ax, 7.62, 1.91, 8.02, 1.75)
    _arrow(ax, 7.62, 1.04, 8.02, 1.20)
    _arrow(ax, 10.42, 1.62, 10.72, 1.62)

    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor=SURFACE)
    plt.close(fig)


def fig_pipeline(path: Path) -> None:
    """The research workflow, with the verification gate made explicit."""
    fig, ax = plt.subplots(figsize=(13.0, 3.5), dpi=200)
    ax.set_xlim(0, 13)
    ax.set_ylim(0, 3.5)
    ax.axis("off")

    steps = [
        ("Raw EAV\ncorpus", "42 subjects", SERIES_2),
        ("Verify\nthe data", "audit + 74 tests", WARN),
        ("Preprocess\nto cache", "~2.8 GB", SERIES_2),
        ("Train\nencoders", "CPU only", SERIES_1),
        ("7-fold\nsubject CV", "every subject\nheld out once", SERIES_1),
        ("Significance\ntesting", "paired, by subject", SERIES_1),
        ("Model of\nrecord", "traced artifact", SERIES_3),
    ]
    w, gap = 1.56, 0.31
    x = 0.15
    for i, (label, sub, colour) in enumerate(steps):
        ax.add_patch(plt.Rectangle((x, 1.15), w, 1.25, facecolor=colour,
                                   edgecolor="none", zorder=2))
        ax.text(x + w / 2, 1.94, label, ha="center", va="center", fontsize=11,
                color=SURFACE, fontweight="bold", zorder=3, linespacing=1.25)
        ax.text(x + w / 2, 0.86, sub, ha="center", va="top", fontsize=9.5,
                color=INK_MUTED, linespacing=1.3)
        if i < len(steps) - 1:
            _arrow(ax, x + w + 0.04, 1.78, x + w + gap - 0.04, 1.78)
        x += w + gap

    verify_cx = 0.15 + (w + gap) + w / 2      # centre of the "Verify the data" step
    ax.annotate("", xy=(verify_cx, 2.48), xytext=(verify_cx, 2.90),
                arrowprops=dict(arrowstyle="-", color=WARN, linewidth=1.6))
    ax.text(verify_cx + 0.22, 2.95, "the step this project originally skipped",
            fontsize=11, color=WARN, ha="left", va="center", fontweight="bold")

    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor=SURFACE)
    plt.close(fig)


def fig_defect_two(path: Path) -> None:
    """Schematic of defect 2: one EEG tensor joined to 100 audio files."""
    fig, axes = plt.subplots(1, 2, figsize=(12.4, 4.5), dpi=200)

    for ax, title, distinct in (
        (axes[0], "Before  —  what the loader built", 1),
        (axes[1], "After  —  the corrected join", 100),
    ):
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 6.4)
        ax.axis("off")
        ax.set_title(title, fontsize=14, color=INK_PRIMARY, pad=14, loc="left")

        for i in range(5):
            y = 5.0 - i * 0.95
            ax.add_patch(plt.Rectangle((0.3, y), 2.3, 0.62, color=SERIES_2, alpha=0.92))
            ax.text(1.45, y + 0.31, f"audio {i + 1}", ha="center", va="center",
                    fontsize=10.5, color=SURFACE)
            eeg_colour = SERIES_1 if distinct == 1 else plt.get_cmap("Blues")(0.45 + i * 0.09)
            ax.add_patch(plt.Rectangle((6.5, y), 2.6, 0.62, color=eeg_colour))
            tag = "EEG  seg[0]" if distinct == 1 else f"EEG  trial {i + 1}"
            ax.text(7.8, y + 0.31, tag, ha="center", va="center", fontsize=10.5,
                    color=SURFACE)
            ax.annotate("", xy=(6.4, y + 0.31), xytext=(2.7, y + 0.31),
                        arrowprops=dict(arrowstyle="->", color=BASELINE, linewidth=1.6))

        ax.text(0.3, 0.55, "⋮   100 audio files", fontsize=11, color=INK_MUTED)
        note = ("all 100 rows share one\nbyte-identical EEG tensor"
                if distinct == 1 else "each row carries its own\nEEG trial")
        ax.text(6.5, 0.35, note, fontsize=11.5,
                color=WARN if distinct == 1 else INK_SECONDARY,
                fontweight="bold" if distinct == 1 else "normal")

    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor=SURFACE)
    plt.close(fig)


# =============================================================================
# Slide construction helpers
# =============================================================================

class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = SW, SH
        self.blank = self.prs.slide_layouts[6]
        self.n = 0
        # (slide index, figure filename), recorded so the build can emit a
        # manifest -- the deck's images are recoverable by hand if a transfer
        # or upload damages the embedded copies.
        self.figure_log: list[tuple[int, str]] = []

    # -- primitives ----------------------------------------------------------
    def _box(self, slide, left, top, width, height):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        return tf

    def _run(self, para, text, *, size, bold=False, colour=INK_PRIMARY, font=FONT,
             italic=False):
        r = para.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = rgb(colour)
        return r

    # -- slide kinds ---------------------------------------------------------
    def slide(self, kicker: str | None, title: str, *, subtitle: str | None = None):
        s = self.prs.slides.add_slide(self.blank)
        self.n += 1
        top = MARGIN - Inches(0.15)

        if kicker:
            tf = self._box(s, MARGIN, top, CONTENT_W, Inches(0.3))
            p = tf.paragraphs[0]
            self._run(p, kicker.upper(), size=12, bold=True, colour=ACCENT)
            top += Inches(0.42)

        tf = self._box(s, MARGIN, top, CONTENT_W, Inches(0.75))
        p = tf.paragraphs[0]
        self._run(p, title, size=31, bold=True, colour=INK_PRIMARY)

        # accent rule under the title
        line = s.shapes.add_shape(1, MARGIN, top + Inches(0.72), Inches(1.05), Pt(3.5))
        line.fill.solid()
        line.fill.fore_color.rgb = rgb(ACCENT)
        line.line.fill.background()
        line.shadow.inherit = False

        s._content_top = top + Inches(1.0)

        if subtitle:
            tf = self._box(s, MARGIN, s._content_top, CONTENT_W, Inches(0.4))
            p = tf.paragraphs[0]
            self._run(p, subtitle, size=15, colour=INK_SECONDARY)
            s._content_top += Inches(0.62)

        self._footer(s)
        return s

    def section(self, number: str, title: str, blurb: str):
        s = self.prs.slides.add_slide(self.blank)
        self.n += 1
        bar = s.shapes.add_shape(1, Emu(0), Emu(0), Inches(0.28), SH)
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(ACCENT)
        bar.line.fill.background()
        bar.shadow.inherit = False

        tf = self._box(s, Inches(1.4), Inches(2.55), Inches(10.5), Inches(2.4))
        p = tf.paragraphs[0]
        self._run(p, number, size=13, bold=True, colour=ACCENT)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(10)
        self._run(p2, title, size=40, bold=True, colour=INK_PRIMARY)
        p3 = tf.add_paragraph()
        p3.space_before = Pt(16)
        self._run(p3, blurb, size=17, colour=INK_SECONDARY)
        self._footer(s)
        return s

    def title_slide(self, title: str, subtitle: str, meta: list[str]):
        s = self.prs.slides.add_slide(self.blank)
        self.n += 1
        bar = s.shapes.add_shape(1, Emu(0), Emu(0), SW, Inches(0.22))
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(ACCENT)
        bar.line.fill.background()
        bar.shadow.inherit = False

        tf = self._box(s, MARGIN, Inches(2.15), Inches(11.3), Inches(2.6))
        p = tf.paragraphs[0]
        self._run(p, title, size=42, bold=True, colour=INK_PRIMARY)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(18)
        self._run(p2, subtitle, size=19, colour=INK_SECONDARY)

        tf2 = self._box(s, MARGIN, Inches(5.5), Inches(11.3), Inches(1.4))
        for i, line in enumerate(meta):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p.space_before = Pt(5)
            self._run(p, line, size=13.5, colour=INK_MUTED)
        return s

    def _footer(self, s):
        tf = self._box(s, MARGIN, SH - Inches(0.62), CONTENT_W, Inches(0.3))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        self._run(p, str(self.n - 1), size=11, colour=INK_MUTED)

    # -- content blocks ------------------------------------------------------
    def bullets(self, s, items, *, left=None, width=None, top=None, size=17,
                gap=13, height=None):
        left = MARGIN if left is None else left
        width = CONTENT_W if width is None else width
        top = s._content_top if top is None else top
        tf = self._box(s, left, top, width, height or (SH - top - Inches(0.8)))
        for i, item in enumerate(items):
            if isinstance(item, str):
                item = {"text": item}
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(0 if i == 0 else gap)
            level = item.get("level", 0)
            p.level = level
            marker = item.get("marker", "•  " if level == 0 else "–  ")
            colour = item.get("colour", INK_PRIMARY if level == 0 else INK_SECONDARY)
            sz = item.get("size", size if level == 0 else size - 2.5)
            if level:
                p.space_before = Pt(6)
            self._run(p, marker, size=sz, colour=item.get("marker_colour", ACCENT),
                      bold=True)
            for seg in _segments(item["text"]):
                self._run(p, seg.text, size=sz, bold=seg.bold, colour=seg.colour or colour,
                          font=seg.font or FONT)
        return tf

    def paragraph(self, s, text, *, top=None, size=17, colour=INK_SECONDARY,
                  left=None, width=None, height=None):
        left = MARGIN if left is None else left
        width = CONTENT_W if width is None else width
        top = s._content_top if top is None else top
        tf = self._box(s, left, top, width, height or Inches(0.55))
        p = tf.paragraphs[0]
        for seg in _segments(text):
            self._run(p, seg.text, size=size, bold=seg.bold, colour=seg.colour or colour,
                      font=seg.font or FONT)
        return tf

    def explain(self, s, text, *, top=None, size=14.5, left=None, width=None):
        """A plain-language gloss for the jargon on the slide.

        Set in muted italic under the technical content, so a reader who already
        knows the term can skip the line without it competing for attention.
        """
        left = MARGIN if left is None else left
        width = CONTENT_W if width is None else width
        top = s._content_top if top is None else top
        tf = self._box(s, left, top, width, Inches(0.8))
        p = tf.paragraphs[0]
        self._run(p, "In plain terms   ", size=size - 1.5, bold=True, colour=ACCENT)
        for seg in _segments(text):
            r = self._run(p, seg.text, size=size, bold=seg.bold,
                          colour=seg.colour or INK_SECONDARY, font=seg.font or FONT)
            r.font.italic = not seg.font
        return tf

    def callout(self, s, text, *, top=None, tone=ACCENT, size=17, left=None,
                width=None, height=Inches(0.95)):
        left = MARGIN if left is None else left
        width = CONTENT_W if width is None else width
        top = (SH - Inches(1.75)) if top is None else top
        band = s.shapes.add_shape(1, left, top, width, height)
        band.fill.solid()
        band.fill.fore_color.rgb = rgb("f4f6f9" if tone == ACCENT else "fbf2ee")
        band.line.fill.background()
        band.shadow.inherit = False
        edge = s.shapes.add_shape(1, left, top, Inches(0.06), height)
        edge.fill.solid()
        edge.fill.fore_color.rgb = rgb(tone)
        edge.line.fill.background()
        edge.shadow.inherit = False

        tf = band.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.25)
        p = tf.paragraphs[0]
        for seg in _segments(text):
            self._run(p, seg.text, size=size, bold=seg.bold,
                      colour=seg.colour or INK_PRIMARY, font=seg.font or FONT)
        return band

    def table(self, s, headers, rows, *, top=None, widths=None, size=14,
              col_align=None, height=None, left=None, total_width=None):
        top = s._content_top if top is None else top
        left = MARGIN if left is None else left
        total_width = CONTENT_W if total_width is None else total_width
        nrows, ncols = len(rows) + 1, len(headers)

        # PowerPoint treats row height as a minimum and grows it to fit, which
        # silently pushes anything placed below. Size each row for its tallest
        # cell up front so the caller's `top` values stay meaningful.
        def row_height(cells) -> Emu:
            lines = max(str(c).count("\n") + 1 for c in cells)
            return Inches(0.34 + 0.26 * (lines - 1) + 0.006 * (size - 14))

        heights = [row_height(headers)] + [row_height(r) for r in rows]
        height = height or Emu(sum(int(h) for h in heights))
        shape = s.shapes.add_table(nrows, ncols, left, top, total_width, height)
        tbl = shape.table
        tbl.first_row = True

        if widths:
            scale = total_width / sum(widths)
            for i, w in enumerate(widths):
                tbl.columns[i].width = Emu(int(w * scale))

        tbl.rows[0].height = heights[0]
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb("eef1f5")
            _cell_text(self, cell, h, size=size - 1, bold=True, colour=INK_SECONDARY,
                       align=(col_align or ["l"] * ncols)[j])

        for i, row in enumerate(rows, start=1):
            tbl.rows[i].height = heights[i]
            for j, val in enumerate(row):
                cell = tbl.cell(i, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(SURFACE if i % 2 else "fafafa")
                _cell_text(self, cell, val, size=size, colour=INK_PRIMARY,
                           align=(col_align or ["l"] * ncols)[j])
        return tbl

    def picture(self, s, name, *, top=None, max_h=None, max_w=None):
        path = FIG_DIR / name
        if not path.exists():
            path = REPORT_FIGS / name
        if not path.exists():
            print(f"  ! missing figure: {name}", file=sys.stderr)
            return None
        top = s._content_top if top is None else top
        max_h = max_h or (SH - top - Inches(1.0))
        max_w = max_w or CONTENT_W

        from PIL import Image
        with Image.open(path) as im:
            iw, ih = im.size
        scale = min(max_w / iw, max_h / ih)
        w, h = Emu(int(iw * scale)), Emu(int(ih * scale))
        self.figure_log.append((self.n - 1, name))
        return s.shapes.add_picture(str(path), Emu(int((SW - w) / 2)), top,
                                    width=w, height=h)

    def code(self, s, lines, *, top=None, size=13.5, left=None, width=None,
             height=None):
        top = s._content_top if top is None else top
        left = MARGIN if left is None else left
        width = CONTENT_W if width is None else width
        height = height or Inches(0.3 + 0.28 * len(lines))
        panel = s.shapes.add_shape(1, left, top, width, height)
        panel.fill.solid()
        panel.fill.fore_color.rgb = rgb("f6f6f4")
        panel.line.fill.background()
        panel.shadow.inherit = False
        tf = panel.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.28)
        tf.margin_top = Inches(0.14)
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(2)
            for seg in _segments(line):
                self._run(p, seg.text, size=size, bold=seg.bold,
                          colour=seg.colour or INK_SECONDARY, font=MONO)
        return panel

    def notes(self, s, text: str) -> None:
        s.notes_slide.notes_text_frame.text = text.strip()

    def save(self, path: Path) -> None:
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))


class _Seg:
    __slots__ = ("text", "bold", "colour", "font")

    def __init__(self, text, bold=False, colour=None, font=None):
        self.text, self.bold, self.colour, self.font = text, bold, colour, font


def _segments(text: str):
    """Split on **bold**, ``mono`` and {red|...} markers."""
    out, buf, i = [], "", 0
    while i < len(text):
        if text.startswith("**", i):
            j = text.find("**", i + 2)
            if j > 0:
                if buf:
                    out.append(_Seg(buf)); buf = ""
                out.append(_Seg(text[i + 2:j], bold=True))
                i = j + 2
                continue
        if text.startswith("``", i):
            j = text.find("``", i + 2)
            if j > 0:
                if buf:
                    out.append(_Seg(buf)); buf = ""
                out.append(_Seg(text[i + 2:j], font=MONO))
                i = j + 2
                continue
        if text.startswith("{red|", i):
            j = text.find("}", i)
            if j > 0:
                if buf:
                    out.append(_Seg(buf)); buf = ""
                out.append(_Seg(text[i + 5:j], bold=True, colour=WARN))
                i = j + 1
                continue
        buf += text[i]
        i += 1
    if buf:
        out.append(_Seg(buf))
    return out or [_Seg("")]


def _cell_text(deck, cell, text, *, size, bold=False, colour=INK_PRIMARY, align="l"):
    tf = cell.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Inches(0.12)
    cell.margin_top = cell.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
    for seg in _segments(str(text)):
        deck._run(p, seg.text, size=size, bold=bold or seg.bold,
                  colour=seg.colour or colour, font=seg.font or FONT)


# =============================================================================
# The deck
# =============================================================================

def build() -> Deck:
    """The 23-slide defence deck.

    Budgeted at 22:20 against a 20--25 minute allowance, so roughly 2.5 minutes
    of headroom remain at the cap. Every slide's notes open with its own budget
    and the cumulative time it should end at, so the presenter can recover pace
    at any point without doing arithmetic on stage.
    """
    d = Deck()

    # =========================================================================
    # 1. TITLE                                                    0:20 -> 0:20
    # =========================================================================
    s = d.title_slide(
        "Multimodal Emotion Recognition with\nCross-Modal Attention Fusion",
        "Recognising emotion from brain signals and speech, and measuring "
        "honestly whether combining them helps",
        [
            "Tuan Tran   ·   Master of Software Engineering",
            "Supervisor:  [SUPERVISOR NAME]",
            "FPT University   ·   [FACULTY]",
            "Capstone defence   ·   [DEFENCE DATE]",
        ],
    )
    d.notes(s, """
⏱  0:20   ·   should end at  0:20

ALLOWANCE: 40 min total -- 20-25 presentation + demo, then 15 minutes of
questions. This deck is budgeted at 22:20, leaving ~2.5 min of headroom.

Every slide's notes start with its budget and the clock time it should end at.
If you are more than 45 seconds behind at slide 12 (RQ1), use the recovery plan
in slides/README.md rather than speeding up.

FILL IN before presenting: supervisor, faculty, defence date.

Opening line: "This project asks a simple question -- if a person's brain activity
and their voice both carry emotional information, does combining the two work
better than either alone? The answer turned out to depend entirely on HOW you
combine them."
""")

    # =========================================================================
    # 2. BACKGROUND & MOTIVATION                                  1:05 -> 1:25
    # =========================================================================
    s = d.slide("2 · Background & motivation", "Two very different windows on one emotion")
    d.table(s,
            ["", "Speech", "EEG (brain activity)"],
            [
                ["What it is", "The sound of the voice — pitch,\nloudness, rhythm",
                 "Electrical activity measured by\nsensors on the scalp"],
                ["Strength", "Cheap; needs only a microphone",
                 "Not deliberately performed —\nharder to fake"],
                ["Weakness", "People hide feelings; noise\ndestroys the signal",
                 "Expensive hardware; varies\nenormously between people"],
            ],
            widths=[1.7, 4.15, 4.15], size=13.5, top=Inches(2.15))
    d.explain(s, "EEG records the tiny electrical signals your brain produces, through 30 "
                 "sensors resting on the scalp — 500 readings per second, per sensor.",
              top=Inches(4.6))
    d.callout(s, "The two fail in **different situations**. That is the whole premise for "
                 "combining them: where one is weak, the other may still be right.",
              top=Inches(5.5), height=Inches(0.95))
    d.notes(s, """
⏱  1:05   ·   should end at  1:25

Do not spend time motivating affective computing in general -- the committee knows
why it matters. One sentence: systems that interact with people work better when
they can tell how the person feels.

The point that must land is the last line, because it is the premise the whole
thesis tests, and section 9 eventually confirms it with numbers.

If asked about facial video: present in the corpus, the original implementation's
video path was a stub returning nothing, excluded from scope rather than reported
as contributing. Named as the most obvious extension.
""")

    # =========================================================================
    # 3. RESEARCH PROBLEM / GAP                                   1:05 -> 2:30
    # =========================================================================
    s = d.slide("3 · Research problem & gap", "Why combining them is harder than it sounds")
    d.bullets(s, [
        "**Different clocks and different quantities.** EEG arrives 500×/second across 30 "
        "sensors; audio features every 10 ms. One measures activity in the head, the other "
        "sound in the air — they share no natural scale",
        "**Few participants.** Corpora with synchronised EEG and audio usually have fewer "
        "than 100 people",
        "**People differ enormously.** A model that learns one person's brain patterns often "
        "fails on the next person",
    ], top=Inches(2.15), size=16, gap=10, height=Inches(2.2))
    d.table(s,
            ["Existing fusion approach", "Limitation"],
            [
                ["**Early** — glue features together, then classify",
                 "Loses each signal's structure"],
                ["**Late** — classify separately, combine answers",
                 "Learns no interaction between the signals"],
                ["**Gated** — learn how much of each to let through",
                 "No explicit model of how the two relate"],
            ],
            widths=[4.6, 5.4], size=13, top=Inches(4.5))
    d.callout(s, "**The gap:** cross-modal attention was proposed to fix this — but it is "
                 "rarely tested on **people the model has never seen**.",
              top=Inches(6.15), height=Inches(0.8), size=16)
    d.notes(s, """
⏱  1:05   ·   should end at  2:30

Three difficulties, one sentence each. The third is the one that explains every
result later, so give it the emphasis.

The table is the standard taxonomy -- point at it, do not read it.

The callout is the gap statement. It has two halves: an untested mechanism and a
testing protocol that is too easy. Most published work trains and tests on the
same participants, which measures something much easier than what a deployed
system faces.
""")

    # =========================================================================
    # 4. OBJECTIVES                                               0:50 -> 3:20
    # =========================================================================
    s = d.slide("4 · Objectives", "What this thesis set out to do")
    d.paragraph(s, "**General objective** — build and honestly evaluate a system that "
                   "recognises emotion by combining EEG and speech through cross-modal "
                   "attention.", top=Inches(2.15), size=16.5, colour=INK_PRIMARY,
                height=Inches(0.75))
    d.table(s,
            ["", "Research question"],
            [
                ["RQ1", "Does EEG carry emotion information that transfers to new people?"],
                ["RQ2", "Does cross-modal attention fusion beat the stronger single signal?"],
                ["RQ3", "How much does an easier evaluation protocol inflate the score?"],
                ["RQ4", "What verification does a multimodal pipeline require?"],
            ],
            widths=[0.6, 9.4], size=14, top=Inches(3.1))
    d.explain(s, "Each results slide is labelled with the research question it answers, so "
                 "you can always see which one is being addressed.", top=Inches(5.5))
    d.notes(s, """
⏱  0:50   ·   should end at  3:20

Read the four RQs aloud. They are the spine of the talk.

Specific objectives, if asked: build a verified pipeline for EAV, design the
fusion architecture, measure it against each signal alone on unseen people, and
establish which comparisons are statistically trustworthy.

A fifth question from the original proposal, on generative augmentation, was
descoped by argument rather than experiment. Mention only if asked.
""")

    # =========================================================================
    # 5. LITERATURE REVIEW                                        0:55 -> 4:15
    # =========================================================================
    s = d.slide("5 · Related work", "What the field reports, and what is different here")
    d.table(s,
            ["Strategy", "Reported behaviour in the literature"],
            [
                ["Feature-level (early) fusion", "Usually beats late fusion by 5–15%"],
                ["Attention-based fusion", "Reported 2–8% over plain concatenation"],
                ["Multi-head attention (Vaswani et al., 2017)",
                 "The mechanism behind modern language models; applied\nacross modalities, each signal learns what to attend to"],
            ],
            widths=[3.9, 6.1], size=13, top=Inches(2.15))
    d.callout(s, "This ranking is the **received wisdom this thesis set out to apply** — and "
                 "section 9 reports the opposite ordering on unseen people.",
              top=Inches(4.35), height=Inches(0.85), size=16)
    d.table(s,
            ["This thesis differs by", ""],
            [
                ["Evaluation", "Every result on **people the model never saw**, cross-validated"],
                ["Verification", "The data pipeline is **proven correct**, not assumed"],
                ["Reporting", "Paired significance tests, and a **measured noise floor**"],
            ],
            widths=[2.1, 7.9], size=13, top=Inches(5.45))
    d.notes(s, """
⏱  0:55   ·   should end at  4:15

Foreshadow deliberately: the literature ordering is early > late, and this thesis
finds late > early on unseen people. Planting it here makes section 9 land as a
result rather than as a surprise.

Attribute fairly -- those published margins are mostly measured under pooled
protocols, which is a comparability problem, not a claim that the authors erred.

If anyone asks whether "+4.25%" is relative or absolute, the answer is
ABSOLUTE: 62.77% went to 67.02%, an arithmetic difference of 4.25. It is not
4.25% OF 62.77. Have that sentence ready -- it is the one ambiguity in the
notation this deck uses.
""")

    # =========================================================================
    # 6. METHODOLOGY                                              0:50 -> 5:05
    # =========================================================================
    s = d.slide("6 · Methodology", "How the study was run")
    d.picture(s, "pipeline.png", top=Inches(2.15), max_h=Inches(2.6))
    d.bullets(s, [
        "**Quantitative and experimental.** One variable changes at a time — encoders, "
        "loss, optimiser, schedule and data partition are identical across every comparison",
        "Training reads a **preprocessed cache**, so every run sees byte-identical inputs",
    ], top=Inches(5.05), size=15.5, gap=10)
    d.notes(s, """
⏱  0:50   ·   should end at  5:05

Walk the diagram left to right in one sentence.

Emphasise the second box: verification comes BEFORE training in this workflow.
That ordering is the methodological lesson of the whole project, and slide 11
shows what happened when it was skipped.
""")

    s = d.slide("6 · Methodology", "The data, and how it reaches the model")
    d.table(s,
            ["", ""],
            [
                ["Corpus", "**EAV** — 42 participants, synchronised EEG, audio and video"],
                ["Used here", "The 100 *Speaking* trials per person carrying matched audio\n"
                              "→ **4,200 multimodal samples**"],
                ["Emotions", "Neutral, Anger, Calmness, Sadness, Happiness — **exactly balanced**"],
                ["Chance level", "**20.0%** — read every accuracy in this talk against this"],
            ],
            widths=[1.7, 8.3], size=13.5, top=Inches(2.15))
    d.table(s,
            ["Each signal is", "EEG", "Speech"],
            [
                ["Cleaned", "Noise and drift filtered out", "Standardised to one sample rate"],
                ["Summarised", "How much energy in each\nbrain rhythm", "A picture of the sound over time"],
                ["Put on a common footing", "Each person aligned to a\nshared reference", "Each frequency band standardised"],
            ],
            widths=[2.1, 3.95, 3.95], size=12.5, top=Inches(4.35))
    d.explain(s, "Neither signal goes into the model raw — each is first turned into a "
                 "compact summary that keeps the emotional content and drops the rest.",
              top=Inches(6.3))
    d.notes(s, """
⏱  1:00   ·   should end at  6:05

SKIM the bottom table -- one sentence: "both signals are cleaned, summarised, and
put on a common footing before the model sees them." The committee does not need
the filter cut-offs, and the details are in the methodology chapter.

Technical detail only if asked: EEG band-pass 0.5-45 Hz then decimated to 125 Hz,
band-power differential entropy per channel per second, Euclidean alignment per
subject. Audio at 16 kHz, 64-band log-mel with a 16 ms hop, SpecAugment in
training.

Two details worth saying aloud, because they are where corpora usually go wrong:
labels come from the corpus's own per-trial label matrix rather than being guessed
from filenames, and the classes are exactly balanced at 840 samples each, so no
hidden imbalance correction is doing work.

The technical point to be ready to defend: the per-person alignment is computed
from that person's OWN recordings and uses no labels. Applying it to a held-out
person corresponds to a brief unlabelled calibration recording, which any real
deployment would have. No label information crosses the split.

Also applied to audio during training: SpecAugment, which hides random stripes of
the spectrogram so the model cannot memorise one recording. Mention only if the
slide is going quickly.
""")

    # =========================================================================
    # 7. PROPOSED SOLUTION / SYSTEM DESIGN                        1:25 -> 7:30
    # =========================================================================
    s = d.slide("7 · System design", "The architecture, and the two ways to combine")
    d.picture(s, "architecture.png", top=Inches(1.95), max_h=Inches(4.15))
    d.callout(s, "Both paths use the **identical encoders**. The only thing that differs is "
                 "how the two streams come together — which is exactly the variable this "
                 "thesis measures.", top=Inches(6.2), height=Inches(0.8), size=15.5)
    d.notes(s, """
⏱  1:25   ·   should end at  7:30

The centrepiece of the design section. Walk it slowly -- this is the slide the
committee will look at longest.

Left to right: each signal goes through its own encoder, which compresses a trial
into 128 numbers -- a compact summary. Both encoders end at the same size, which
is what makes any fusion possible at all. Then the paths diverge: Path A is the
mechanism this thesis proposed, Path B is the trivial alternative.

Say the parameter counts out loud -- 8,064 against 0. That contrast is what the
audience should carry into the results.

Encoder sizes if asked: EEG 190,897 parameters, audio 444,352. Each can also run
alone into the same classifier, which is how the single-signal comparisons are
kept fair.
""")

    s = d.slide("7 · System design", "Cross-modal attention fusion — the proposed mechanism")
    d.bullets(s, [
        "**Each signal looks at the other, in both directions.** The brain signal asks “what "
        "in the voice is relevant to me?”, and the voice asks the same of the brain signal",
        "**Each signal also keeps its own view**, so if one is uninformative the system can "
        "fall back on the other",
        "**A learned switch** then decides how much of each to pass on to the final decision",
    ], top=Inches(2.2), size=16.5, gap=12, height=Inches(2.5))
    d.explain(s, "“Attention” means the model learns *which parts* of one signal to look at "
                 "when interpreting the other, rather than treating everything equally.",
              top=Inches(4.85))
    d.callout(s, "**8,064 parameters** in this module — every one of them learned from the "
                 "training participants.", top=Inches(5.75), height=Inches(0.85), size=17)
    d.notes(s, """
⏱  1:05   ·   should end at  8:35

Keep this conceptual -- the three bullets as written are the whole explanation.
Do NOT put the equations on screen; they are on backup slide B1 and only come out
if someone asks.

This is your own contribution to the architecture. Present it with conviction --
that it loses is the results section's business, not this slide's.

Formal detail is on backup slide B1 if asked: A' = MHA(E, A, A) and
E' = MHA(A, E, E), each with residual add and layer norm, then a sigmoid gate over
the concatenated pair. 4 heads, 32 dims per head.

The final phrase -- "learned from the training participants" -- is the seed of the
explanation for why it fails. Say it deliberately and let it sit.
""")

    s = d.slide("7 · System design", "Four ways to combine, one identical protocol")
    d.table(s,
            ["Mechanism", "Fusion params", "Idea"],
            [
                ["Cross-modal attention (pooled)", "8,064", "Attend between the two trial summaries"],
                ["Cross-modal attention (sequence)", "8,064", "Attend before summarising, to align events in time"],
                ["Subject-adversarial fusion", "8,064", "Also punish the model for identifying *who* the person is"],
                ["**Late fusion — mean of probabilities**", "**0**", "**Average the two models' answers**"],
            ],
            widths=[4.0, 1.7, 4.3], size=13.5, top=Inches(2.2))
    d.explain(s, "Each of the first three targets a different suspected cause of failure: "
                 "too little capacity, lost timing information, and person-specific patterns.",
              top=Inches(4.85))
    d.callout(s, "Only the combiner changes between rows — which is what makes the "
                 "comparison meaningful.", top=Inches(5.75), height=Inches(0.85), size=16)
    d.notes(s, """
⏱  0:50   ·   should end at  9:25

The experimental design in one table, and what turns a single negative result into
a systematic one: three principled attempts, each aimed at a different hypothesis,
plus the trivial baseline.

Late fusion was added LAST, after the diagnosis on slide 14 pointed at it. Say
that -- it shows the ordering was driven by evidence, not by trying everything.
""")

    # =========================================================================
    # 8. IMPLEMENTATION / EXPERIMENT                              1:00 -> 10:25
    # =========================================================================
    s = d.slide("8 · Implementation & evaluation", "How the system is built and tested",
                subtitle="Every number in this talk is measured on people the model has never seen.")
    d.bullets(s, [
        "**7-fold subject-wise cross-validation** — the 42 participants split into 7 groups; "
        "each group is held out in turn, giving **4,200 predictions** on unseen people",
        "**Python 3.11 + PyTorch, CPU only.** The full experiment matrix runs in **under an "
        "hour**; 74 automated tests; seeds fixed and CPU threads pinned",
        "**Metrics:** accuracy and unweighted average recall, with paired significance tests",
    ], top=Inches(2.65), size=16, gap=11, height=Inches(2.5))
    d.explain(s, "**Subject-independent** means no participant appears in both training and "
                 "testing. It is much harder than the common alternative — and much closer to "
                 "how the system would really be used.", top=Inches(5.2))
    d.callout(s, "**RQ3 answered:** the identical model scores **+12.75% higher** when "
                 "participants may appear on both sides — about **17× the largest "
                 "architectural effect** measured here.",
              top=Inches(6.1), height=Inches(0.9), size=15.5)
    d.notes(s, """
⏱  1:00   ·   should end at  10:25

The most important methodological slide in the deck. If you are behind schedule,
take the time from elsewhere, not from here.

SKIM the second bullet (tooling) in about five seconds -- "it runs on a laptop in
under an hour and has 74 automated tests". Spend the time on the first and third.

Say the number: 68.25% pooled against 55.50% subject-independent. Same weights,
same code, only the partition changed.

Why 7 folds and not leave-one-out: 42/7 = 6 held out per fold, which keeps each
test fold large enough for the by-subject bootstrap, at about four hours rather
than seventeen.

The thread-pinning detail looks like a footnote and is not -- slide 16 shows it is
worth 2.2%.
""")

    s = d.slide("8 · Implementation", "Verifying the data before trusting any result")
    d.table(s,
            ["Defect found in the original pipeline", "Consequence"],
            [
                ["EEG array read along the wrong axis", "Each EEG 'trial' was a single 2-millisecond instant"],
                ["**One EEG recording shared by 100 samples**",
                 "**The EEG input carried no emotion information at all**"],
                ["Splits let the same person train and test", "69% of the test set had been trained on"],
                ["Failed loads silently replaced by noise", "The amount of noise in any result was unknowable"],
            ],
            widths=[4.6, 5.4], size=13, top=Inches(2.2))
    d.callout(s, "All four produced **well-formed data and believable accuracy**. None "
                 "produced an error message. Every result predating the fix is withdrawn.",
              top=Inches(4.85), tone=WARN, height=Inches(0.9), size=16)
    d.bullets(s, [
        "Now checked automatically: labels cross-referenced against the corpus's own label "
        "matrix, all **12,600 media files** verified, failures raise errors instead of "
        "substituting data",
    ], top=Inches(5.95), size=15)
    d.notes(s, """
⏱  0:55   ·   should end at  11:20

Compressed from four slides into one. Do not expand it -- but do not skip it. It
is Contribution 1 and it is why the numbers that follow can be believed.

The second row is the fatal one: each participant has one EEG file and 100 audio
files, and the join paired every audio file with the same EEG recording. So the
EEG branch received identical input carrying five different labels. Everything
previously reported as "multimodal fusion" was in substance audio-only.

How they were found, if asked: while building the verification harness for the
subject-independent split -- not by noticing a suspicious number. The numbers
looked fine. That is the point.
""")

    # =========================================================================
    # 9. RESULTS                                                  4:35 -> 15:55
    # =========================================================================
    s = d.slide("9 · Results  ·  RQ1", "Each signal on its own")
    d.table(s,
            ["Model", "Accuracy on unseen people", "vs chance"],
            [
                ["Random guessing", "20.00%", "—"],
                ["**EEG only**", "**45.32% ± 0.92**", "**+25.3%**"],
                ["Audio only", "62.77% ± 1.52", "+42.8%"],
            ],
            widths=[4.4, 3.4, 2.2], col_align=["l", "c", "c"], size=15, top=Inches(2.25))
    d.bullets(s, [
        "**RQ1 answered: yes.** EEG reaches 2.3× chance on people the model has never seen",
        "The old pipeline was {red|structurally incapable} of showing this — its EEG input "
        "was identical for every sample",
        "But audio is **17.5% stronger**. Any fusion must improve on the better of the "
        "two, not the worse",
    ], top=Inches(4.25), size=16, gap=11)
    d.notes(s, """
⏱  0:45   ·   should end at  12:05

The "±" figures are standard deviations across repeated runs with different random
seeds. Explain that phrase once here and it carries.

The last bullet is the bar everything else has to clear. Say it clearly -- it
prevents the common misreading that fusion only needs to beat EEG.
""")

    s = d.slide("9 · Results  ·  RQ2", "Does cross-modal attention fusion help?")
    d.table(s,
            ["Approach", "Fusion params", "Accuracy", "vs audio alone", "p"],
            [
                ["Audio only", "—", "64.48%", "—", "—"],
                ["**Cross-modal attention (pooled)**", "**8,064**", "**64.12%**", "**−0.36%**",
                 "**0.72** ❌"],
                ["Cross-modal attention (sequence)", "8,064", "63.19%", "−1.29%", "❌"],
                ["Subject-adversarial fusion", "8,064", "63.60%", "−0.88%", "❌"],
            ],
            widths=[4.0, 1.7, 1.5, 1.8, 1.0], col_align=["l", "c", "c", "c", "c"],
            size=13.5, top=Inches(2.2))
    d.explain(s, "**p** is the probability of seeing a difference this large by pure chance. "
                 "Below 0.05 is normally called real; **0.72 means we cannot tell this apart "
                 "from no difference at all**.", top=Inches(4.4))
    d.callout(s, "All three learned mechanisms fail. Per participant it is a coin flip — "
                 "audio wins 23 of 42, fusion 18, one tie. **The mechanism this thesis "
                 "proposed does not work.**", top=Inches(5.4), tone=WARN, height=Inches(0.95),
              size=16)
    d.notes(s, """
⏱  0:50   ·   should end at  12:55

Deliver this without hedging. It is a result, not a failure to report.

Three principled attempts, three nulls -- that is what makes it a systematic
finding rather than one unlucky architecture.

Then pause before advancing, because the next slide overturns the obvious
interpretation.
""")

    s = d.slide("9 · Results", "Is EEG simply redundant? No.",
                subtitle="Instead of reading the average, look at the trials audio gets wrong.")
    d.table(s,
            ["Question", "Answer"],
            [
                ["How many trials does audio get wrong?", "**1,537 of 4,200**"],
                ["On those exact trials, how often is EEG right?", "**47.50%**  (chance is 20%)"],
                ["If we could always pick the better model per trial?", "**80.79%**  —  +17.4% over audio"],
            ],
            widths=[5.6, 4.4], size=14.5, top=Inches(2.6))
    d.callout(s, "**EEG is 2.4× chance precisely where audio fails.** The attention module "
                 "recovers 791 of audio's mistakes but breaks 761 of its correct answers — a "
                 "net gain of **30 trials out of 4,200**. A near one-for-one trade, not an "
                 "absence of signal.", top=Inches(4.9), height=Inches(1.2), size=16)
    d.explain(s, "The flat average was hiding a busy exchange, not an empty channel — the "
                 "information is there, and the combiner was spending it as fast as it gained "
                 "it.", top=Inches(6.25))
    d.notes(s, """
⏱  1:00   ·   should end at  13:55

The intellectual core of the thesis. This is the analysis that turned a dead end
into a direction, so protect its time.

Conditioning on the audio model's errors separates two explanations that demand
opposite responses: "EEG is redundant" (abandon the modality) versus "EEG is
complementary but unexploited" (fix the combiner).

The 791 / 761 pair is the sentence to land.

Per-class detail is on backup slide B4 -- happiness is the one emotion where EEG
beats audio, and there fusion beats both by 24.9%.
""")

    s = d.slide("9 · Results  ·  RQ2", "The combination that does work")
    d.paragraph(s, "Average the two models' answers. That is the entire method — no fusion "
                   "parameters at all.", top=Inches(2.2), size=17, colour=INK_PRIMARY,
                height=Inches(0.6))
    d.table(s,
            ["System", "Fusion params", "Accuracy", "vs audio alone", "p"],
            [
                ["Audio only", "—", "62.77% ± 1.52", "—", "—"],
                ["Cross-modal attention fusion", "8,064", "64.12%", "−0.36%", "0.72 ❌"],
                ["**Late fusion — mean of probabilities**", "**0**", "**67.02% ± 1.11**",
                 "**+4.25% ± 0.41**", "**0.003** ✅"],
                ["*Perfect per-trial choice (unreachable)*", "—", "*81.69%*", "*+17.21%*", "—"],
            ],
            widths=[4.2, 1.7, 1.9, 1.9, 1.0], col_align=["l", "c", "c", "c", "c"],
            size=13.5, top=Inches(3.0))
    d.callout(s, "Positive on **every repetition**, worst case +3.79%, and better on "
                 "**32 of the 42 participants** — a broad effect, not a few outliers.",
              top=Inches(5.7), height=Inches(0.9), size=16.5)
    d.notes(s, """
⏱  1:05   ·   should end at  15:00

The headline. Slow right down.

Say it plainly: "the best system in this project has no fusion parameters at all,
and the 8,064-parameter module I designed is beaten by an average."

If challenged on the overlapping error bars -- 67.02 +- 1.11 against 62.77 +- 1.52
-- the two move together, because changing the seed redraws the same fold
assignment for both models. The paired difference has a spread of only +-0.41,
which is why the difference rather than the accuracy is the quantity quoted.
""")

    s = d.slide("9 · Results", "The pattern behind the result, and how small a difference is real")
    d.picture(s, "combiner_freedom.png", top=Inches(1.95), max_h=Inches(3.3))
    d.callout(s, "**The fewer parameters the combiner has, the better it transfers to new "
                 "people** — it learns which signal to trust *for the training participants*, "
                 "and that preference does not carry across.",
              top=Inches(5.4), height=Inches(0.85), size=15.5)
    d.bullets(s, [
        "Measured noise floor: **2.2%** from CPU thread scheduling alone, ±0.92% from the "
        "seed. **Anything under ~1% is not reproducible** — the +4.25% gain sits 3.9× "
        "outside it",
    ], top=Inches(6.4), size=14.5)
    d.notes(s, """
⏱  0:55   ·   should end at  15:55

SKIM the chart mechanics. Say the pattern in one sentence -- "the fewer parameters
the combiner has, the better it transfers" -- then the noise floor in one more, and
move on. The detail below is for questions, not for the talk.

Two things on one slide because they defend each other: the pattern, and the
yardstick that says the pattern is real.

Why paired differences rather than accuracies: the two runs' audio baselines
differ by 2.19% at the same configuration, inside the noise floor, so absolute
accuracies across runs are not comparable.

What replicates across both runs, and IS the claim: the equal average is the best
rule in both; the 8,064-parameter module is the worst in both. What does NOT
replicate -- concede immediately if asked -- is the order of the 1-parameter and
5-parameter rules, which swaps between runs at under 0.5% apart.

On the noise floor: running the same code on a different number of CPU threads
changes the answer, because the machine adds numbers in a different order. Nothing
about the model changed. A single-seed figure of 69.55% this project reported on
9 August is withdrawn on the same standard.
""")

    # =========================================================================
    # DEMO                                                        3:00 -> 18:55
    # =========================================================================
    s = d.slide("Demonstration", "The system running on people it has never seen")
    d.code(s, [
        "$ python scripts/demo_replay.py --auto --delay 2",
        "",
        "  trial 07/20   truth: Happiness",
        "    multimodal   -> Happiness   0.61    [====================      ]",
        "    audio only   -> Anger       0.44    [==============            ]",
    ], top=Inches(2.2))
    d.bullets(s, [
        "Replays trials from the **8 held-out participants** — subjects 5, 19, 22, 26, 27, "
        "30, 40 and 42, excluded from training and validation entirely",
        "Each trial shows the true emotion, both models' predictions, and their confidence",
        "Replaying all 800 of their trials reproduces the recorded test accuracy **exactly** "
        "— so the demo path and the evaluation path are provably the same code",
    ], top=Inches(4.2), size=15.5, gap=11)
    d.notes(s, """
⏱  3:00   ·   should end at  18:55

BEFORE THE SESSION: have the terminal already open, virtual environment activated,
in the project root, with the command typed but not run. Font size up. Never
install or debug on stage.

FALLBACK: if anything fails, a pre-recorded screen capture of the same run is the
correct response -- do not troubleshoot in front of the committee. Record it in
advance regardless; you may never need it.

WHAT TO NARRATE (about 8-10 trials, ~2 minutes):
  - point out a trial where both models agree and are right
  - point out a Happiness trial where the multimodal model is right and audio is
    wrong -- this is the complementarity result made visible
  - point out one where both are wrong, and say so plainly: at 67% roughly one
    trial in three is wrong, and the demo does not hide that

CLOSING LINE for the demo: "these are eight people the model has never seen, and
the numbers on screen are the same numbers in the results table."

If asked why no live microphone: EAV participants are Korean speakers reading
scripted prompts in studio conditions. A presenter speaking English into a laptop
mic is a different distribution with no ground truth, so a wrong prediction would
be uninterpretable rather than informative. That refusal is deliberate and
documented.
""")

    # =========================================================================
    # 10. DISCUSSION                                              1:30 -> 20:25
    # =========================================================================
    s = d.slide("10 · Discussion", "Why the simple method wins")
    d.bullets(s, [
        "**The hard problem here is new people, not new sounds.** One person's brain "
        "patterns differ from the next person's more than one emotion differs from another",
        "Every parameter in a learned combiner is fitted on the **training participants** — "
        "and what it fits is *which signal to trust*, which is exactly what varies by person",
        "The plain average is the only combiner that **cannot** overfit that, because it has "
        "nothing to fit",
    ], top=Inches(2.2), size=16.5, gap=12, height=Inches(2.8))
    d.callout(s, "Direct evidence: a single weight tuned on the validation participants "
                 "scored **69.00%** on them and **64.88%** on new ones — *worse* than not "
                 "tuning it at all. **One parameter was enough to overfit six people.**",
              top=Inches(5.25), height=Inches(1.05), size=16)
    d.notes(s, """
⏱  0:50   ·   should end at  19:45

The explanation slide. If the committee accepts this mechanism, the whole result
hangs together.

The callout is the cleanest single piece of evidence for the mechanism, and it
does not depend on any ranking of the rules.

Supporting evidence if pressed: the EEG model's validation accuracy exceeds its
test accuracy by 6.42% while audio shows the reverse -- even between two groups
of unseen people, the EEG representation is unstable.
""")

    s = d.slide("10 · Discussion", "Strengths, limits, and what remains open")
    d.table(s,
            ["", ""],
            [
                ["**Strengths**", "Verified pipeline · every result on unseen people · four\n"
                                  "mechanisms under one protocol · a measured noise floor"],
                ["**Limitations**", "One corpus, 42 participants · scripted Korean speech in studio\n"
                                    "conditions · three seeds · the null results were not repeated"],
                ["**Not implemented**", "Video modality · generative augmentation (descoped by\n"
                                        "argument) · artefact removal"],
                ["**Still open**", "The perfect-choice ceiling is 80.79% — about **12%** of\n"
                                   "complementary information remain unexploited"],
            ],
            widths=[1.9, 8.1], size=13.5, top=Inches(2.2))
    d.notes(s, """
⏱  0:40   ·   should end at  20:25

Volunteering limitations before the committee finds them is worth more than an
extra result.

Say plainly that this is a research artifact, not a deployable system: 67% on five
balanced classes is not accurate enough where a mistake carries cost.

The last row is the bridge to future work.
""")

    # =========================================================================
    # 11. CONTRIBUTIONS                                           0:45 -> 21:10
    # =========================================================================
    s = d.slide("11 · Contributions", "What this thesis adds")
    d.bullets(s, [
        "**Academic — a bounded, diagnosed negative result.** Cross-modal attention fusion "
        "does not beat audio alone on unseen people (−0.36%, p = 0.72), with the cause "
        "localised to transfer through a *learned* combiner",
        "**Academic — a positive result from a zero-parameter combiner.** Averaging beats "
        "audio by +4.25% ± 0.41 (p = 0.003); combiner freedom, not fusion capacity, is the "
        "operative variable",
        "**Methodological — a verified pipeline.** Four silent defects documented with "
        "evidence, and the four cheap practices that catch them",
        "**Practical — a reproducible artifact.** The full experiment matrix runs in under "
        "an hour on CPU, every number traced to a recorded file",
    ], top=Inches(2.2), size=16, gap=12)
    d.notes(s, """
⏱  0:45   ·   should end at  21:10

Four contributions, in the thesis's own order. Do not add a fifth.

If asked "what is new, averaging is not new": the contribution is not the
mechanism, it is the measurement that identifies WHEN the trivial mechanism is
right, and the diagnosis that explains why.
""")

    # =========================================================================
    # 12. CONCLUSION & FUTURE WORK                                0:50 -> 22:00
    # =========================================================================
    s = d.slide("12 · Conclusion & future work", "Answers, and where this goes next")
    d.table(s,
            ["", "Question", "Answer"],
            [
                ["RQ1", "Does EEG transfer to new people?",
                 "**Yes** — 45.32% vs 20% chance, and 47.5% on\nexactly the trials audio gets wrong"],
                ["RQ2", "Does attention fusion beat the\nstronger signal?",
                 "**No** (−0.36%, p = 0.72) — but averaging the\ntwo models does: **+4.25%**, p = 0.003"],
                ["RQ3", "How much does an easier protocol\ninflate the score?", "**+12.75%**"],
                ["RQ4", "What verification is required?", "Four silent defects, four practices, 74 tests"],
            ],
            widths=[0.6, 4.1, 5.3], size=12.5, top=Inches(2.15))
    d.bullets(s, [
        "**Next, in priority order:** per-trial confidence gating (the 80.79% ceiling says "
        "~12% are available) → subject normalisation of EEG → add the video modality",
    ], top=Inches(5.35), size=15)
    d.callout(s, "The architecture was the problem — **not the modality**.", top=Inches(6.15),
              height=Inches(0.75), size=18)
    d.notes(s, """
⏱  0:50   ·   should end at  22:00

Read the answers column only; the questions are there for the committee to read.

Future work is ranked by how likely each is to change the conclusion, and each
targets a measured failure rather than a hypothesised one -- the same discipline
as the rest of the thesis.

The closing line is the thesis in six words. Pause after it.
""")

    # =========================================================================
    # 13. Q&A                                                     0:20 -> 22:20
    # =========================================================================
    s = d.slide(None, "Thank you", subtitle="Questions")
    d.bullets(s, [
        "**EEG carries real emotion signal** — strongest exactly where speech fails",
        "**Cross-modal attention fusion does not beat speech alone** on people the model has "
        "never seen",
        "**Averaging the two models does** — +4.25% ± 0.41, with zero fusion parameters",
    ], top=Inches(2.9), size=17, gap=14)
    d.paragraph(s, "Every number traced to docs/CHANGELOG.md   ·   "
                   "github.com/ttuan8600/MSE-CAPSTONE-PROJECT",
                top=Inches(5.6), size=14, colour=INK_MUTED)
    d.notes(s, """
⏱  0:20   ·   should end at  22:20

Three sentences, in this order, then stop talking. Do not add a fourth.

15 minutes of questions follow. The prepared answers are in slides/QA_PREP.md --
read it the morning of the defence, not the night before.

Backup slides: B1 attention equations · B2 defect evidence · B3 single partition
vs CV · B4 per-class exchange · B5 confusion matrix · B6 statistics ·
B7 adversarial retraction · B8 deployment · B9 reproduction commands.
""")

    # =========================================================================
    # BACKUP
    # =========================================================================
    s = d.section("Backup", "Backup slides", "Not presented — for the 15 minutes of questions.")
    d.notes(s, "Do not advance past this in the talk itself.")

    s = d.slide("Backup 1", "Cross-modal attention — formal detail")
    d.code(s, [
        "A'  = MultiHeadAttention(E, A, A)      # EEG queries audio",
        "E'  = MultiHeadAttention(A, E, E)      # audio queries EEG",
        "",
        "Ê   = LayerNorm(E' + E)                # residual, per stream",
        "Â   = LayerNorm(A' + A)",
        "",
        "g   = sigmoid(W_g · [Ê ; Â])           # learned gate",
        "F   = LayerNorm(MLP([ g⊙Ê ; (1-g)⊙Â ]))",
        "",
        "4 heads × 32 dims = 128     MLP: Linear(256,128) → ReLU → Drop(0.1) → Linear(128,128)",
    ], top=Inches(2.3))
    d.bullets(s, [
        "Fusion module **8,064** · classifier 66,565 · full fusion model 850,417 parameters",
    ], top=Inches(5.6), size=15)
    d.notes(s, "Residual connections are why the module can degrade gracefully towards a "
               "single modality -- and it does exactly that.")

    s = d.slide("Backup 2", "Defect evidence, verbatim")
    d.code(s, [
        "$ python scripts/verify_data_fix.py",
        "",
        " subject  samples   original  corrected",
        "       1      100          1        100",
        "       2      100          1        100",
        "       3      100          1        100",
        "",
        "Across 5 subjects (500 samples): 5 distinct EEG tensors originally,",
        "                                 500 after the fix.",
        "EEG tensor shape: [28, 200] (original) -> [30, 2500] (corrected)",
    ], top=Inches(2.3))
    d.notes(s, "The array ships as (10000, 30, 200) = (time, channels, trials). The loader "
               "read eeg[0,:,:] believing it was the first segment; it is time index 0 -- a "
               "single 2 ms sample.")

    s = d.slide("Backup 3", "Why one held-out group was not enough")
    d.picture(s, "cross_validation.png", top=Inches(2.1), max_h=Inches(3.6))
    d.callout(s, "A single 8-participant partition overstated audio by 3.73% and understated "
                 "fusion by 1.50% — **reversing the ranking**. Per-participant accuracy spans "
                 "25% to 81%.", top=Inches(5.95), height=Inches(0.9), size=15.5)
    d.notes(s, "Orange diamonds are the single partition; bars are 7-fold CV. This is the "
               "empirical justification for the extra compute.")

    s = d.slide("Backup 4", "Where the exchange happens, per emotion")
    d.picture(s, "complementarity.png", top=Inches(2.05), max_h=Inches(3.7))
    d.callout(s, "**Happiness** is the one emotion where EEG beats audio — and there fusion "
                 "beats both by **+24.9%**. On the other four, audio is already strong and "
                 "fusion costs 4–7%.", top=Inches(5.95), height=Inches(0.9), size=15.5)
    d.notes(s, "Happiness is the weakest vocal cue in this corpus, and the neural channel "
               "compensates for exactly that weakness. These recalls are for the attention "
               "module, so they are the anatomy of the null result.")

    s = d.slide("Backup 5", "Confusion matrix — model of record")
    d.picture(s, "confusion_matrix.png", top=Inches(2.1), max_h=Inches(4.3))
    d.notes(s, "Characteristic confusions are Calmness with Neutral and Happiness with Anger "
               "-- pairs sharing prosodic contours, which is what the EEG channel is best "
               "placed to disambiguate.")

    s = d.slide("Backup 6", "Statistical method")
    d.bullets(s, [
        "**Paired McNemar exact test** — both models score the identical 4,200 samples",
        "**Bootstrap confidence intervals**, 10,000 resamples, resampled **by participant** — "
        "the unit of independence is the person, not the trial",
        "**Multi-seed aggregation** — the seed changes both the fold assignment and the "
        "initialisation, so the spread describes the whole procedure",
        "The paired difference (±0.41) is far more stable than the individual accuracies "
        "(±1.11, ±1.52), because pairing cancels the shared fold draw",
    ], top=Inches(2.3), size=16, gap=13)
    d.notes(s, "Resampling by trial would treat 100 trials from one person as 100 independent "
               "observations and badly understate the interval.")

    s = d.slide("Backup 7", "Subject-adversarial training: retracted")
    d.table(s,
            ["Claim", "Seed 42 only", "Across seeds", "p", "Verdict"],
            [
                ["Adversarial EEG, alone", "+0.64%", "+0.18% ± 0.82 (n=4)", "0.695", "❌"],
                ["Adversarial EEG, in fusion", "+1.29%", "+0.92% ± 0.61 (n=3)", "0.121", "❌"],
            ],
            widths=[3.4, 1.8, 2.5, 1.2, 1.1], col_align=["l", "c", "c", "c", "c"],
            size=14, top=Inches(2.4))
    d.bullets(s, [
        "Per seed, alone: −0.91, +0.02, +0.95, +0.65% — **the sign flips**",
        "Kept because the negative result is informative: forcing the EEG representation to "
        "be person-independent by gradient reversal is **not** what unlocks the gain. The "
        "**combination rule** is",
    ], top=Inches(4.0), size=16, gap=12)
    d.notes(s, "The cleanest illustration in the project of why one seed is not evidence.")

    s = d.slide("Backup 8", "Deployment")
    d.table(s,
            ["", ""],
            [
                ["Model of record", "Audio-only — 510,917 parameters, 1.95 MB, CPU inference"],
                ["Why audio-only", "Chosen when fusion was measured as tied; 3.4× more stable\n"
                                   "across folds and needs no EEG cap"],
                ["Where EEG is available", "Late fusion is the better system (+4.25%) and is what a\n"
                                           "deployment should adopt"],
                ["API", "`POST /predict` returns emotion + full distribution;\n"
                        "a missing modality returns HTTP 400, never silent zeros"],
            ],
            widths=[2.3, 7.7], size=13, top=Inches(2.2))
    d.notes(s, "The deployable-model choice and the scientific result are different "
               "decisions, and Chapter 6 splits the recommendation by deployment context "
               "rather than naming one winner.")

    s = d.slide("Backup 9", "Reproducing every number in this talk")
    d.code(s, [
        "pytest                                        # 74 passed",
        "",
        "python scripts/audit_eav_alignment.py         # verify the corpus",
        "python scripts/preprocess_eav.py              # ~5 min, writes ~2.8 GB",
        "",
        "python scripts/cross_validate_late_fusion.py  # the headline result",
        "python scripts/aggregate_seeds.py             # multi-seed aggregation",
        "python scripts/analyze_complementarity.py     # the diagnosis",
        "python scripts/compare_cv.py outputs/cv_*     # McNemar + bootstrap",
        "python scripts/generate_result_figures.py     # figures, from artifacts only",
    ], top=Inches(2.3))
    d.callout(s, "The complete experiment matrix runs in **under an hour on CPU** from a "
                 "verified cache.", top=Inches(5.8), height=Inches(0.8), size=16)
    d.notes(s, "Offer to run any of these live.")

    return d


# Figures the deck reuses from the report rather than rendering itself. They are
# copied into slides/figures/ so that directory holds every image the deck uses
# -- one folder to hand over, and one place to look when re-inserting an image by
# hand after a damaged transfer.
REUSED_FIGURES = (
    "complementarity.png",
    "cross_validation.png",
    "confusion_matrix.png",
)


def write_manifest(deck: Deck, path: Path) -> None:
    """Record which figure sits on which slide, for manual recovery."""
    used = sorted(set(deck.figure_log))
    spare = sorted(
        p.name for p in FIG_DIR.glob("*.png")
        if p.name not in {name for _, name in used}
    )
    lines = [
        "# Figures used by MSE_Capstone_Defence.pptx",
        "",
        "Regenerate everything with `python scripts/make_slides.py`.",
        "Slide numbers are 0-based, matching the number printed on each slide.",
        "",
        "| Slide | Figure | Section |",
        "| --- | --- | --- |",
    ]
    section = {
        "pipeline.png": "6 · Methodology",
        "architecture.png": "7 · System design",
        "combiner_freedom.png": "9 · Results",
        "cross_validation.png": "Backup 3",
        "complementarity.png": "Backup 4",
        "confusion_matrix.png": "Backup 5",
    }
    for idx, name in used:
        lines.append(f"| {idx} | `{name}` | {section.get(name, '—')} |")
    if spare:
        lines += [
            "",
            "## Not used by the current deck",
            "",
            *(f"- `{n}`" for n in spare),
            "",
            "Kept because they are cheap to regenerate and useful when answering "
            "questions.",
        ]
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def main() -> int:
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    print("Rendering deck figures ...")
    fig_architecture(FIG_DIR / "architecture.png")
    fig_pipeline(FIG_DIR / "pipeline.png")
    fig_combiner_freedom(FIG_DIR / "combiner_freedom.png")
    fig_multiseed(FIG_DIR / "multiseed.png")
    fig_defect_two(FIG_DIR / "defect_two.png")

    print("Copying reused report figures ...")
    for name in REUSED_FIGURES:
        source = REPORT_FIGS / name
        if not source.exists():
            print(f"  ! missing source figure: {source}", file=sys.stderr)
            continue
        shutil.copyfile(source, FIG_DIR / name)
        print(f"  {name}")

    print("Building slides ...")
    deck = build()
    out = SLIDE_DIR / "MSE_Capstone_Defence.pptx"
    deck.save(out)
    write_manifest(deck, FIG_DIR / "MANIFEST.md")
    print(f"Wrote {out.relative_to(ROOT)}  ({deck.n} slides, "
          f"{len(set(deck.figure_log))} figures)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
